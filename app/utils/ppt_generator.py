from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
import os
import json
from openai import OpenAI
import requests
import uuid
import re
from io import BytesIO
from typing import List, Dict, Optional

class PPTGenerator:
    def __init__(self):
        api_key = os.environ.get('OPENAI_API_KEY')
        if not api_key:
            raise ValueError("OpenAI API key not found in environment variables. Please set OPENAI_API_KEY in your .env file.")
        
        print(f"Debug - API Key loaded in PPTGenerator: {api_key[:10]}...")
        
        try:
            self.client = OpenAI(api_key=api_key)
            print("Debug - OpenAI client initialized successfully")
            
            # Test the API key with a simple completion request
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "system", "content": "Test message"}],
                max_tokens=5
            )
            print("Debug - OpenAI API key verified successfully")
            
        except Exception as e:
            print(f"Debug - Error with OpenAI setup: {str(e)}")
            if 'Invalid API key' in str(e):
                raise ValueError("Invalid OpenAI API key. Please check your .env file.")
            elif 'Rate limit' in str(e):
                raise ValueError("OpenAI API rate limit reached. Please try again later.")
            else:
                raise ValueError(f"OpenAI API error: {str(e)}")

        # Define available template styles
        self.TEMPLATE_STYLES = {
            "Aesthetic": "Aesthetic.pptx",


            "Vintage": "Vintage.pptx",

            "Creative": "Creative.pptx",

            "Simple": "Simple.pptx",

            "Clean and Neat": "Clean and Neat.pptx",
            "Neat and Clean": "Clean and Neat.pptx",
            "Business": "Business.pptx",
            "Verdant": "Verdant.pptx"
        }

        # Define default font styles (fallback)
        self.FONTS = {
            'title': {
                'name': 'Calibri',
                'size': Pt(44),
                'bold': True,
                'color': RGBColor(33, 33, 33)
            },
            'subtitle': {
                'name': 'Calibri',
                'size': Pt(32),
                'bold': False,
                'color': RGBColor(0, 123, 255)
            },
            'body': {
                'name': 'Calibri',
                'size': Pt(18),
                'bold': False,
                'color': RGBColor(33, 37, 41)
            }
        }
        
    # Image generation helper
    def _generate_image(self, prompt: str) -> str:
        """Generate an image using DALL·E 3 and save it locally. Returns the file path or empty string on failure."""
        try:
            print(f"Debug - Generating image for prompt: {prompt}")
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                n=1,
                size="1024x1024"
            )
            image_url = response.data[0].url
            img_bytes = requests.get(image_url).content
            images_dir = os.path.abspath(os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'generated', 'images'))
            os.makedirs(images_dir, exist_ok=True)
            file_path = os.path.join(images_dir, f"{uuid.uuid4().hex}.png")
            with open(file_path, "wb") as f:
                f.write(img_bytes)
            return file_path
        except Exception as e:
            print(f"Warning - Image generation failed: {str(e)}")
            return ""

        # Define available template styles
        self.TEMPLATE_STYLES = {
            "Aesthetic": "Aesthetic.pptx",


            "Vintage": "Vintage.pptx",

            "Creative": "Creative.pptx",

            "Simple": "Simple.pptx",

            "Clean and Neat": "Clean and Neat.pptx",
            "Neat and Clean": "Clean and Neat.pptx",
            "Business": "Business.pptx",
            "Verdant": "Verdant.pptx"
        }
        
        # Define font settings (fallback for missing template styles)
        self.FONTS = {
            'title': {
                'name': 'Calibri',
                'size': Pt(44),
                'bold': True,
                'color': RGBColor(33, 33, 33)
            },
            'subtitle': {
                'name': 'Calibri',
                'size': Pt(32),
                'bold': False,
                'color': RGBColor(0, 123, 255)
            },
            'body': {
                'name': 'Calibri',
                'size': Pt(18),
                'bold': False,
                'color': RGBColor(33, 37, 41)
            }
        }

    def get_template_path(self, style: str) -> str:
        """Get the path to the selected template style"""
        filename = self.TEMPLATE_STYLES.get(style, "Aesthetic.pptx")  # Default to Professional if style not found
        # Get absolute path to the project root directory
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(__file__)))
        template_path = os.path.join(project_root, "app", "static", "presentations", "custom_styles", filename)
        print(f"Debug - Using template path: {template_path}")
        return template_path

    def _remove_all_slides(self, prs: Presentation) -> None:
        """Remove all existing slides from the presentation while preserving the template"""
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)  # Create a list from iterator
        for slide_id in slides:
            xml_slides.remove(slide_id)

    def _apply_text_style(self, shape, style_type='body'):
        """Apply text style to a shape (only used if template styles are missing)"""
        if not shape.has_text_frame:
            return

        style = self.FONTS[style_type]
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = style['name']
            paragraph.font.size = style['size']
            paragraph.font.bold = style['bold']
            paragraph.font.color.rgb = style['color']
            if style_type in ['title', 'subtitle']:
                paragraph.alignment = PP_ALIGN.CENTER

    def _get_layout_by_name(self, prs: Presentation, layout_name: str) -> Optional[any]:
        """Find a slide layout by its name"""
        for layout in prs.slide_layouts:
            if layout.name.lower() == layout_name.lower():
                return layout
        return None

    def _get_title_layout(self, prs: Presentation) -> any:
        """Get the title slide layout
        Strictly selects a proper title slide layout for the first slide.
        """
        # First check if we have any layouts
        if not prs.slide_layouts:
            raise ValueError("Template has no slide layouts")

        # First priority: Look for layout named "Title Slide"
        for layout in prs.slide_layouts:
            if layout.name and "Title Slide" in layout.name:
                print(f"Debug - Found Title Slide layout: {layout.name}")
                return layout

        # Second priority: Check common title slide names
        title_names = ["Title", "Cover", "Cover Page", "Opening"]
        for layout in prs.slide_layouts:
            if any(name.lower() in layout.name.lower() for name in title_names):
                print(f"Debug - Found title layout by name: {layout.name}")
                return layout

        # Last resort: Use the first layout (usually the title layout)
        print(f"Debug - Using first layout as title: {prs.slide_layouts[0].name}")
        return prs.slide_layouts[0]

    def _get_content_layout(self, prs: Presentation) -> any:
        """Get the content slide layout
        Strictly selects a proper content slide layout (not the title slide layout).
        """
        # First check if we have any layouts
        if not prs.slide_layouts:
            raise ValueError("Template has no slide layouts")

        # Skip the title slide layout
        title_layout = self._get_title_layout(prs)

        # First priority: Look for "Title and Content" layout
        for layout in prs.slide_layouts:
            if layout != title_layout and layout.name and "Title and Content" in layout.name:
                print(f"Debug - Found Title and Content layout: {layout.name}")
                return layout

        # Second priority: Look for any content-specific layout
        content_names = ["Content", "Text and Content", "Text", "Title and Text", "Section Header", "Two Content"]
        for layout in prs.slide_layouts:
            if layout != title_layout and any(name.lower() in layout.name.lower() for name in content_names):
                print(f"Debug - Found content layout by name: {layout.name}")
                return layout

        # Last resort: Use any layout that's not the title layout
        for layout in prs.slide_layouts:
            if layout != title_layout:
                print(f"Debug - Using alternate layout: {layout.name}")
                return layout

        # Absolute fallback: Use the second layout if available
        if len(prs.slide_layouts) > 1:
            print(f"Debug - Using second layout: {prs.slide_layouts[1].name}")
            return prs.slide_layouts[1]

        # If all else fails, use any non-first layout
        for i, layout in enumerate(prs.slide_layouts):
            if i > 0:
                print(f"Debug - Using layout {i}: {layout.name}")
                return layout

        print("Warning: Could not find distinct content layout")
        return prs.slide_layouts[0]
        # Content names are already defined above
        for name in layout_names:
            layout = self._get_layout_by_name(prs, name)
            if layout:
                print(f"Debug - Found named content layout: {name}")
                return layout
        
        # Try to find any layout with at least a title placeholder
        for layout in prs.slide_layouts:
            if has_placeholder_type(layout, 1):  # Has title
                print(f"Debug - Using layout with title: {layout.name}")
                return layout
        
        # Final fallback to first layout
        print(f"Debug - Using fallback first layout: {prs.slide_layouts[0].name}")
        return prs.slide_layouts[0]

    def _add_title_slide(self, prs: Presentation, title: str, presenter: str):
        """Add a robust title slide.
        Works even when the selected template has no proper title/subtitle placeholders.
        Removes leftover empty placeholders to avoid slides that just show default prompt text.
        """
        layout = self._get_title_layout(prs)
        slide = prs.slides.add_slide(layout)

        # -------------------------
        # Title handling
        # -------------------------
        title_shape = None
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_frame = title_shape.text_frame
            title_frame.text = title
            title_frame.word_wrap = True
        else:
            # No title placeholder – create our own centred textbox
            width = prs.slide_width * 0.8
            left = int((prs.slide_width - width) / 2)
            top = Inches(1.0)
            height = Inches(1.5)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_shape.text_frame
            title_frame.clear()
            title_frame.text = title
            title_frame.word_wrap = True

        # Dynamic font size based on title length
        length = len(title)
        font_size = 44 if length <= 30 else 40 if length <= 50 else 32
        for para in title_frame.paragraphs:
            para.font.size = Pt(font_size)

        # -------------------------
        # Subtitle handling
        # -------------------------
        subtitle_shape = None
        for shape in slide.placeholders:
            # Standard subtitle placeholder is type 2
            if shape.placeholder_format.type == 2 and shape != title_shape:
                subtitle_shape = shape
                break

        if subtitle_shape is None:
            # Create textbox just below the title
            left = title_shape.left
            top = title_shape.top + title_shape.height + Inches(0.2)
            width = title_shape.width
            height = Inches(1.0)
            subtitle_shape = slide.shapes.add_textbox(left, top, width, height)

        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.clear()
        subtitle_frame.text = f"Presented by {presenter}"
        for para in subtitle_frame.paragraphs:
            para.font.size = Pt(24)

        # -------------------------
        # Clean up stray placeholders
        # -------------------------
        for shp in list(slide.shapes):
            try:
                if not getattr(shp, "is_placeholder", False):
                    continue
                # Keep populated shapes
                if shp in (title_shape, subtitle_shape):
                    continue
                text_ok = False
                if shp.has_text_frame:
                    txt = "".join(p.text for p in shp.text_frame.paragraphs).strip()
                    if txt and not txt.lower().startswith(("click to add", "select this paragraph")):
                        text_ok = True
                if not text_ok:
                    shp._element.getparent().remove(shp._element)
            except Exception:
                # If anything goes wrong, skip deletion to avoid crashing
                continue

    def _add_content_slide(self, prs: Presentation, title: str, content: str):
        """Add content slide"""
        layout = self._get_content_layout(prs)
        print(f"Debug - Using layout: {layout.name} for content slide")
        slide = prs.slides.add_slide(layout)

        # Index of the slide (0-based)
        slide_index = len(prs.slides) - 1

        # ------------------------------------------------------------------
        # Visual enhancement for text-only (free tier) slides: accent side bar
        # ------------------------------------------------------------------
        try:
            sidebar_width = Inches(0.25)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sidebar_width, prs.slide_height)
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(0, 102, 204)  # brand blue
            bar.fill.fore_color.transparency = 0.15  # subtle
            bar.line.fill.background()
        except Exception as e:
            print(f"Debug - Could not add sidebar: {e}")

        # ------------------------------------------------------------------
        # Horizontal rule under title for visual separation
        # ------------------------------------------------------------------
        if slide.shapes.title is not None:
            try:
                rule_top = slide.shapes.title.top + slide.shapes.title.height + Pt(4)
                rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), rule_top,
                                               prs.slide_width - Inches(1.6), Pt(1.5))
                rule.fill.solid()
                rule.fill.fore_color.rgb = RGBColor(0, 0, 0)
                rule.fill.fore_color.transparency = 0.85
                rule.line.fill.background()
            except Exception as e:
                print(f"Debug - Could not add rule: {e}")

        print(f"Debug - Available placeholders in slide: {[f'{ph.placeholder_format.type}:{ph.placeholder_format.idx}' for ph in slide.placeholders]}")
        
        # Add title if title placeholder exists
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Find content placeholder using multiple approaches
        content_placeholder = None
        print("Debug - Searching for content placeholder...")
        for shape in slide.placeholders:
            print(f"Debug - Checking placeholder: type={shape.placeholder_format.type}, idx={shape.placeholder_format.idx}")
            # Skip the title placeholder entirely
            if shape == slide.shapes.title:
                continue

            # Accept this placeholder if it is a body/content placeholder OR an empty text placeholder
            is_body_placeholder = shape.placeholder_format.type in [2, 7]  # Body (2) or Content (7)
            is_common_content_idx = shape.placeholder_format.idx in [1, 2]
            is_empty_text_placeholder = hasattr(shape, 'text') and not shape.text

            if is_body_placeholder or (is_common_content_idx and not is_body_placeholder) or is_empty_text_placeholder:
                content_placeholder = shape
                print(f"Debug - Found content placeholder: type={shape.placeholder_format.type}, idx={shape.placeholder_format.idx}")
                break
        
        # Add content if placeholder exists
        if content_placeholder:
            try:
                # Ensure content placeholder sits below the title to avoid overlap
                if slide.shapes.title is not None:
                    title_bottom = slide.shapes.title.top + slide.shapes.title.height
                    margin = Pt(10)
                    if content_placeholder.top < title_bottom + margin:
                        # Move content placeholder just below title
                        delta = (title_bottom + margin) - content_placeholder.top
                        content_placeholder.top += delta
                        # Reduce height so it still fits on the slide
                        if content_placeholder.height > delta:
                            content_placeholder.height -= delta
                
                # Determine whether to use the placeholder or create a new textbox
                used_placeholder = False  # track if we actually populate the placeholder
                use_placeholder = False
                if content_placeholder is not None:
                    try:
                        # Check text direction attribute on <a:bodyPr>
                        bodyPr = content_placeholder._element.bodyPr
                        vert_attr = bodyPr.get('vert') if bodyPr is not None else None
                        # Use placeholder only if horizontal text orientation and reasonably wide
                        if vert_attr in (None, 'horz') and content_placeholder.width >= Inches(3.2):
                            use_placeholder = True
                    except Exception:
                        pass
                
                if not use_placeholder:
                    # Create our own horizontal textbox on left side
                    left_margin = Inches(0.8)
                    top_margin = Inches(1.5)
                    box_width = prs.slide_width * 0.45  # ~45% of slide width
                    box_height = prs.slide_height - top_margin - Inches(1.0)
                    textbox = slide.shapes.add_textbox(left_margin, top_margin, box_width, box_height)
                    text_frame = textbox.text_frame
                    text_frame.clear()  # Remove default empty paragraph
                    used_placeholder = False
                else:
                    # Use the placeholder's text frame
                    text_frame = content_placeholder.text_frame
                    used_placeholder = True

                # Common text frame settings
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                # ------------------------------------------------------------------
                # Build bullet points with stylish icons & optional two-column layout
                # ------------------------------------------------------------------
                points = [pt.strip() for pt in content.split('\n') if pt.strip()]
                icons = ['▸', '‣', '✓', '✦']
                # Use one bullet style per slide, alternating across slides
                bullet_icon = icons[slide_index % len(icons)]

                def populate_frame(frame, pts, align_right=False):
                    frame.clear()
                    frame.word_wrap = True
                    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    first_local = True
                    for i, txt in enumerate(pts):
                        if first_local:
                            p = frame.paragraphs[0]
                            first_local = False
                        else:
                            p = frame.add_paragraph()
                        # Split main bullet and explainer (first hyphen/en-dash with optional spaces)
                        parts = re.split(r"\s*[-–]\s+", txt, maxsplit=1)
                        if len(parts) == 2:
                            main_txt, explainer = parts
                        else:
                            main_txt, explainer = txt, None

                        p.text = f"{bullet_icon} {main_txt.strip()}"
                        p.level = 0
                        p.space_before = Pt(12)
                        p.space_after = Pt(12)
                        p.line_spacing = 1.2
                        p.alignment = PP_ALIGN.RIGHT if align_right else PP_ALIGN.LEFT

                        # Add explainer as secondary paragraph at level 1, smaller font
                        if explainer:
                            exp_p = frame.add_paragraph()
                            exp_p.text = explainer.strip()
                            exp_p.level = 1
                            exp_p.space_before = Pt(4)
                            exp_p.space_after = Pt(8)
                            exp_p.line_spacing = 1.0
                            exp_p.alignment = PP_ALIGN.RIGHT if align_right else PP_ALIGN.LEFT
                            try:
                                exp_p.font.size = Pt(14)
                                exp_p.font.italic = True
                            except Exception:
                                pass

                # Determine alignment based on slide index (alternate)
                align_right = (slide_index % 2 == 1)

                # Always use two-column layout for better readability
                mid = (len(points) + 1) // 2
                left_pts = points[:mid]
                right_pts = points[mid:]

                col_width = prs.slide_width * 0.4
                col_height = prs.slide_height - Inches(2)
                top_margin = Inches(1.5)
                left_x = Inches(0.8)
                right_x = prs.slide_width - col_width - Inches(0.8)

                left_box = slide.shapes.add_textbox(left_x, top_margin, col_width, col_height)
                right_box = slide.shapes.add_textbox(right_x, top_margin, col_width, col_height)

                populate_frame(left_box.text_frame, left_pts, align_right=False)
                populate_frame(right_box.text_frame, right_pts, align_right=True if align_right else False)

                # Remove any placeholder/textbox created earlier if not needed
                if not use_placeholder:
                    try:
                        slide.shapes._spTree.remove(textbox._element)
                    except Exception:
                        pass
                print("Debug - Successfully added content to slide")

                # Aggressively remove ALL unused placeholders (empty text) except the ones we filled.
                for shp in list(slide.shapes):
                    try:
                        if not shp.is_placeholder:
                            continue
                    except AttributeError:
                        continue

                    # Always keep the title placeholder. Only keep the content placeholder if we actually used it.
                    if shp == slide.shapes.title:
                        continue
                    if shp == content_placeholder and used_placeholder:
                        continue

                    # Determine if the placeholder is effectively empty or just shows the default prompt
                    empty = True
                    if shp.has_text_frame:
                        txt = "".join(p.text for p in shp.text_frame.paragraphs).strip()
                        default_prompt = txt.lower().startswith("click to add")
                        if txt and not default_prompt:
                            empty = False
                    # Treat placeholders with vertical orientation (non-horizontal text) as removable noise
                    try:
                        bodyPr = shp._element.bodyPr
                        vert_attr = bodyPr.get("vert") if bodyPr is not None else None
                        if vert_attr and vert_attr != "horz":
                            empty = True
                    except Exception:
                        pass
                    # If it's a picture or other non-text placeholder with no meaningful content, it's safe to drop.
                    if empty:
                        slide.shapes._spTree.remove(shp._element)
            except Exception as e:
                print(f"Debug - Error adding content to placeholder: {str(e)}")
        else:
            error_msg = "No suitable content placeholder found in the selected template"
            print(f"Debug - {error_msg}")
            raise ValueError(error_msg)


    def generate_title(self, description: str) -> str:
        """Generate an intelligent, professional title from the user's description"""
        try:
            system_prompt = """You are a professional presentation title generator. Your task is to create a polished, 
            engaging title from a given description. The title should be:
            - Professional and presentation-appropriate
            - Concise (max 60 characters)
            - Capture the key theme
            - Use proper title case
            - Optionally use a colon to separate main title and subtitle
            DO NOT use quotes or any special formatting. Return ONLY the title text."""

            user_prompt = f"Create a professional presentation title from this description: {description}"

            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_tokens=60,
                temperature=0.7
            )

            title = response.choices[0].message.content.strip()
            print(f"Debug - Generated title: {title}")
            return title

        except Exception as e:
            print(f"Warning: Could not generate title: {str(e)}. Using description as fallback.")
            return description

    def generate_slide_content(self, prompt: str, num_slides: int) -> List[Dict]:
        """Generate slide content using GPT-3.5"""
        try:
            messages = [
                {
                    "role": "system",
                    "content": (
                        "You are a presentation content generator. Generate a JSON object with exactly this structure:\n"
                        "{\"slides\": [{\"title\": \"string\", \"content\": [\"string\"]}]}"
                        "\nThe 'slides' array MUST contain exactly the requested number of slides where:"
                        "\n- 'title' is the slide title"
                        "\n- 'content' is an array of 4-5 strings where each string is a concise bullet point followed by a brief explanatory sentence (max 20 words) separated by a dash (-)"
                        ""
"\n- The FIRST slide must be an 'Agenda' slide outlining the main sections."
"\n- The LAST slide must be an 'Outro' or 'Conclusion' slide summarizing key takeaways."
                        "\nDo not include any explanation or other text, just the JSON."
                        "\nEnsure you generate exactly the requested number of unique slides."
                    )
                },
                {
                    "role": "user",
                    "content": (
                      f"Create exactly {num_slides} unique slides about: {prompt}. "
                      f"The slides must follow this order: "
                      f"1) Agenda slide; "
                      f"(n-1) topic slides; "
                      f"last slide titled 'Conclusion' or 'Outro'. "
                      f"Each slide must have a unique title and exactly 4-5 bullet points, each followed by a brief explanatory sentence (max 20 words)."
                )
                }
            ]
            
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=messages,
                temperature=0.7,
                max_tokens=2000
            )
        
            # Extract and parse the response
            response_text = response.choices[0].message.content
            print(f"Debug - OpenAI Response: {response_text}")
            response_data = json.loads(response_text)
            
            if not isinstance(response_data, dict) or 'slides' not in response_data:
                raise ValueError("Invalid response format from GPT. Expected object with 'slides' array.")
            
            slides_data = response_data['slides']
            print(f"Debug - Number of slides in response: {len(slides_data)}")
            
            if not isinstance(slides_data, list):
                raise ValueError("Invalid 'slides' format. Expected array.")
            
            # Ensure we have the right number of slides and format the content
            if len(slides_data) < num_slides:
                raise ValueError(f"OpenAI returned only {len(slides_data)} slides, but {num_slides} were requested")

            slides = []
            seen_titles = set()

            for slide in slides_data:
                if not isinstance(slide, dict) or 'title' not in slide or 'content' not in slide:
                    raise ValueError("Invalid slide format. Expected object with 'title' and 'content'.")
                
                if not isinstance(slide['content'], list):
                    raise ValueError("Invalid content format. Expected array of strings.")
                
                if slide['title'] in seen_titles:
                    continue  # Skip duplicate titles
                
                seen_titles.add(slide['title'])
                
                # Join content without bullet points - PowerPoint will add them automatically
                formatted_content = "\n".join(point.strip() for point in slide['content'])
                slides.append({
                    "title": slide['title'],
                    "content": formatted_content
                })
                print(f"Debug - Added slide: {slide['title']}")
                
                if len(slides) == num_slides:
                    break

            if len(slides) < num_slides:
                raise ValueError(f"Could not generate enough unique slides. Got {len(slides)}, needed {num_slides}")
                
            return slides
        except json.JSONDecodeError as e:
            raise Exception(f"Error parsing GPT response: {str(e)}")
        except ValueError as e:
            raise Exception(f"Error in response format: {str(e)}")
        except Exception as e:
            raise Exception(f"Error generating content: {str(e)}")

    def create_presentation(self,
                    title: str,
                    presenter: str,
                    slides_content: List[Dict],
                    template_style: str = "Aesthetic",
                    include_images: bool = False) -> str:
        """Create PowerPoint presentation using a selected template style"""
        # Generate an intelligent title from the input description
        presentation_title = self.generate_title(title)
        # Load template
        template_path = self.get_template_path(template_style)
        if not os.path.exists(template_path):
            raise ValueError(f"Template file not found: {template_path}")
        try:
            prs = Presentation(template_path)
            # Remove any existing slides while preserving the template
            self._remove_all_slides(prs)
        except Exception as e:
            print(f"Warning: Could not load template {template_path}. Using blank presentation. Error: {str(e)}")
            prs = Presentation()

        # Add slides
        print(f"Debug - Adding title slide with generated title: {presentation_title}")
        self._add_title_slide(prs, presentation_title, presenter)
        
        print(f"Debug - Number of content slides to add: {len(slides_content)}")
        for slide_content in slides_content:
            # Add content slide first
            print(f"Debug - Adding content slide: {slide_content['title']}")
            self._add_content_slide(prs, slide_content['title'], slide_content['content'])

            # Optionally add an image generated by DALL·E 3
            if include_images:
                try:
                    img_prompt = f"{slide_content['title']} illustrative image"
                    img_path = self._generate_image(img_prompt)
                    if img_path:
                        # Add the picture roughly on the right half of the slide
                        slide = prs.slides[-1]
                        # Try to insert into a dedicated picture placeholder if present
                        pic_placeholder = None
                        for shp in slide.placeholders:
                            try:
                                if shp.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                                    pic_placeholder = shp
                                    break
                            except Exception:
                                pass

                        # Insert the image and capture the resulting shape reference
                        if pic_placeholder:
                            pic_shape = pic_placeholder.insert_picture(img_path)
                        else:
                            # Fallback: place on right half, respecting slide margins
                            pic_width = Inches(4)
                            left = prs.slide_width - pic_width - Inches(0.5)
                            top = Inches(1.0)  # start a bit higher to leave more space for bottom border
                            pic_shape = slide.shapes.add_picture(img_path, left, top, width=pic_width)

                        # ------------------------------------------------------------------
                        # Post-adjustment: ensure the image does NOT overlap the bottom line
                        # ------------------------------------------------------------------
                        bottom_margin = Inches(0.5)
                        max_height = prs.slide_height - bottom_margin - pic_shape.top
                        if pic_shape.height > max_height:
                            ratio = max_height / pic_shape.height
                            pic_shape.height = int(pic_shape.height * ratio)
                            pic_shape.width = int(pic_shape.width * ratio)

                        # Recalculate left bound / width after possible resize
                        left = pic_shape.left
                        pic_width = pic_shape.width

                        # Reduce width of text-containing shapes to avoid overlap
                        available_width = left - Inches(0.3)
                        for shp in slide.shapes:
                            # Skip pictures
                            if shp.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                                continue
                            if hasattr(shp, "text_frame") and shp.text_frame is not None:
                                # Adjust width if current right edge goes beyond image left OR if placeholder is very narrow
                                # Determine available horizontal space for this shape
                                max_width_allowed = available_width - shp.left
                                min_width_needed = Inches(4)

                                if max_width_allowed <= Inches(1):
                                    continue  # No space to change

                                # If shape is wider than allowed, shrink; if narrower than reasonable, grow (if space)
                                if shp.left + shp.width > available_width:
                                    # shrink to fit but not below min reasonable width
                                    shp.width = max(min_width_needed, max_width_allowed)
                                elif shp.width < min_width_needed and max_width_allowed >= min_width_needed:
                                    # expand to a comfortable width
                                    shp.width = min_width_needed
                        print("Debug - Image added to slide")
                except Exception as e:
                    print(f"Warning - Could not add image to slide: {str(e)}")



        # Get the absolute path to the generated directory
        generated_dir = os.path.abspath(os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'generated'))
        
        # Ensure output directory exists
        os.makedirs(generated_dir, exist_ok=True)

        # Clean filename - more restrictive sanitization
        def sanitize_filename(s: str, max_length: int = 50) -> str:
            # Replace any non-alphanumeric chars with underscore
            safe = ''.join(c if c.isalnum() else '_' for c in s)
            # Remove repeated underscores
            while '__' in safe:
                safe = safe.replace('__', '_')
            # Trim to max length and remove trailing underscores
            safe = safe[:max_length].rstrip('_')
            return safe or 'presentation'  # Fallback if empty
        
        filename = f"{sanitize_filename(title)}.pptx"
        output_path = os.path.join(generated_dir, filename)

        prs.save(output_path)
        return output_path
